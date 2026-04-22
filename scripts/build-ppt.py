#!/usr/bin/env python3
"""Renderiza cada slide de um deck Reveal.js em PNG e monta um .pptx.

Uso:
    python3 scripts/build-ppt.py onboarding-sprint-1

Requer:
    pip install playwright python-pptx --break-system-packages
    python3 -m playwright install chromium

Supõe um servidor estático local servindo o diretório `aulas/`.
Se o servidor não estiver no ar, o script sobe um `python -m http.server` na porta 4799.
"""
import asyncio
import os
import socket
import subprocess
import sys
import time
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

ROOT = Path(__file__).resolve().parents[1]
AULAS_DIR = ROOT / "aulas"
EXPORT_DIR = AULAS_DIR / "exports"
EXPORT_DIR.mkdir(exist_ok=True)

PORT = 4799
WIDTH = 1920
HEIGHT = 1200


def port_open(port: int) -> bool:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.settimeout(0.3)
        try:
            s.connect(("127.0.0.1", port))
            return True
        except OSError:
            return False


def ensure_server():
    if port_open(PORT):
        return None
    proc = subprocess.Popen(
        ["python3", "-m", "http.server", str(PORT), "--directory", str(AULAS_DIR)],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    for _ in range(30):
        if port_open(PORT):
            return proc
        time.sleep(0.2)
    raise RuntimeError("Servidor local não subiu")


async def render_all(deck_name: str, out_dir: Path) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": WIDTH, "height": HEIGHT},
            device_scale_factor=1,
        )
        page = await ctx.new_page()

        # descobre a quantidade de slides horizontais
        base = f"http://localhost:{PORT}/{deck_name}.html"
        await page.goto(base, wait_until="networkidle")
        await page.wait_for_timeout(800)
        total = await page.evaluate(
            "() => document.querySelectorAll('.reveal .slides > section').length"
        )

        for i in range(total):
            url = f"{base}#/{i}"
            await page.goto(url, wait_until="networkidle")
            await page.wait_for_timeout(600)
            out = out_dir / f"slide-{i+1:02d}.png"
            await page.screenshot(path=str(out), full_page=False)
            print(f"  slide {i+1:02d} renderizado")

        await browser.close()
        return total


def build_pptx(deck_name: str, slides_dir: Path, total: int) -> Path:
    prs = Presentation()
    # 1920x1200 é 16:10. Inches usa 96dpi de referência.
    prs.slide_width = Emu(int(13.333 * 914400))   # 13.333" largura
    prs.slide_height = Emu(int(8.333 * 914400))   # 8.333" altura

    blank_layout = prs.slide_layouts[6]

    for i in range(total):
        slide = prs.slides.add_slide(blank_layout)
        img_path = slides_dir / f"slide-{i+1:02d}.png"
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    out_path = EXPORT_DIR / f"{deck_name}.pptx"
    prs.save(str(out_path))
    return out_path


def main():
    deck_name = sys.argv[1] if len(sys.argv) > 1 else "onboarding-sprint-1"
    server_proc = ensure_server()
    try:
        frames_dir = EXPORT_DIR / f"{deck_name}-frames"
        total = asyncio.run(render_all(deck_name, frames_dir))
        out = build_pptx(deck_name, frames_dir, total)
        size_kb = out.stat().st_size // 1024
        print(f"\nPPT gerado em {out} ({size_kb} KB, {total} slides)")
    finally:
        if server_proc is not None:
            server_proc.terminate()


if __name__ == "__main__":
    main()
